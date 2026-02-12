import os
import shutil
from pptx import Presentation
from PyPDF2 import PdfReader

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import Chroma


# =====================================================
# LOAD PPT (FAST — NO NLTK)
# =====================================================
def load_ppt(file_path):
    print(f"📂 Loading PPT: {file_path}")

    prs = Presentation(file_path)
    full_text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text += shape.text + "\n"

    return Document(page_content=full_text)


# =====================================================
# LOAD PDF (FAST — BOOK MODE)
# =====================================================
def load_pdf(file_path):
    print(f"📘 Loading PDF: {file_path}")

    reader = PdfReader(file_path)
    full_text = ""

    for page in reader.pages:
        text = page.extract_text()
        if text:
            full_text += text + "\n"

    return Document(page_content=full_text)


# =====================================================
# SPLIT INTO SMART CHUNKS
# =====================================================
def split_documents(documents):
    print("✂️ Splitting into chunks (BOOK MODE)...")

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,      # GOOD FOR BOOKS
        chunk_overlap=200
    )

    chunks = splitter.split_documents(documents)

    print(f"✅ Created {len(chunks)} chunks")
    return chunks


# =====================================================
# CREATE VECTOR DATABASE
# =====================================================
def create_vector_db(chunks):
    print("🧠 Creating embeddings...")

    embedding_model = HuggingFaceEmbeddings(
        model_name="sentence-transformers/all-MiniLM-L6-v2"
    )

    vector_db = Chroma.from_documents(
    documents=chunks,
    embedding=embedding_model,
    persist_directory="./vector_db"
)


    print("✅ Vector DB saved!")
    return vector_db


# =====================================================
# 🔥 BOOK MODE INGESTION (AUTO SCAN DATA FOLDER)
# =====================================================
def run_ingestion(data_folder="data"):

    print("🚀 Starting BOOK MODE ingestion...")

    all_documents = []

    # Scan data folder automatically
    for filename in os.listdir(data_folder):

        file_path = os.path.join(data_folder, filename)

        # PPT
        if filename.endswith(".pptx"):
            doc = load_ppt(file_path)
            all_documents.append(doc)

        # PDF
        elif filename.endswith(".pdf"):
            doc = load_pdf(file_path)
            all_documents.append(doc)

    if not all_documents:
        print("❌ No supported files found in data folder")
        return None

    # Split into chunks
    chunks = split_documents(all_documents)

    # 🔥 Rebuild vector DB every time (important for new books)
    if os.path.exists("vector_db"):
        shutil.rmtree("vector_db")
        print("♻️ Old vector DB deleted")

    vector_db = create_vector_db(chunks)

    print("🎉 BOOK MODE ingestion complete!")
    return vector_db
